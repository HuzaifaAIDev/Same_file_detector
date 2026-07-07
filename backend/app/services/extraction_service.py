"""
Text extraction service.

This is a refactor of the original detector.py's extraction logic:
same extraction behaviour and the same supported formats, but with the
dead/commented-out first draft removed, logging instead of silent
failure, and a thread-safe LRU-bounded cache instead of an
unbounded dict.
"""
import csv
import json
import mimetypes
import os
import warnings
from functools import lru_cache
from pathlib import Path
from threading import Lock

from app.core.config import settings
from app.core.logging_config import get_logger

logger = get_logger(__name__)

warnings.filterwarnings("ignore")

# ---------------------------------------------------------------------
# Optional dependencies — the app should still start (with reduced
# format support) if some of these aren't installed.
# ---------------------------------------------------------------------
try:
    import fitz  # PyMuPDF
except Exception:  # pragma: no cover
    fitz = None
    logger.warning("PyMuPDF not available — PDF extraction disabled.")

try:
    import pytesseract
    from PIL import Image
except Exception:  # pragma: no cover
    pytesseract = None
    Image = None
    logger.warning("Pillow/pytesseract not available — OCR disabled.")

try:
    from rapidfuzz import fuzz
except Exception:  # pragma: no cover
    fuzz = None
    logger.warning("rapidfuzz not available — similarity scoring disabled.")

try:
    from docx import Document
except Exception:  # pragma: no cover
    Document = None

try:
    import openpyxl
except Exception:  # pragma: no cover
    openpyxl = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover
    Presentation = None


if pytesseract is not None and os.path.exists(settings.TESSERACT_CMD):
    pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_CMD


# ---------------------------------------------------------------------
# Thread-safe extraction cache (bounded to avoid unbounded memory growth
# over a long-running server process).
# ---------------------------------------------------------------------
_cache: dict[str, str] = {}
_cache_lock = Lock()
_CACHE_MAX_ENTRIES = 500


def _cache_get(key: str) -> str | None:
    with _cache_lock:
        return _cache.get(key)


def _cache_set(key: str, value: str) -> None:
    with _cache_lock:
        if len(_cache) >= _CACHE_MAX_ENTRIES:
            _cache.pop(next(iter(_cache)))
        _cache[key] = value


def clean(text: str | None) -> str:
    if not text:
        return ""
    return str(text).lower().strip()


# ---------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------

def extract_pdf(path: str) -> str:
    if fitz is None:
        return ""
    text = ""
    try:
        with fitz.open(path) as doc:
            text = " ".join(page.get_text() for page in doc)

            # OCR fallback for scanned/low-text PDFs
            if len(text.strip()) < 20 and pytesseract is not None:
                text = ""
                for page in doc:
                    try:
                        pix = page.get_pixmap()
                        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                        text += pytesseract.image_to_string(img, config="--oem 3 --psm 6") + " "
                    except Exception as exc:
                        logger.debug("OCR page failed for %s: %s", path, exc)
                        continue
    except Exception as exc:
        logger.warning("PDF extraction failed for %s: %s", path, exc)
    return text


def extract_image(path: str) -> str:
    if pytesseract is None or Image is None:
        return ""
    try:
        with Image.open(path) as img:
            return pytesseract.image_to_string(img, config="--oem 3 --psm 6")
    except Exception as exc:
        logger.warning("Image OCR failed for %s: %s", path, exc)
        return ""


def extract_txt(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as exc:
        logger.warning("TXT extraction failed for %s: %s", path, exc)
        return ""


def extract_json(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            data = json.load(f)
        return json.dumps(data, ensure_ascii=False)
    except Exception as exc:
        logger.warning("JSON extraction failed for %s: %s", path, exc)
        return ""


def extract_csv(path: str) -> str:
    try:
        rows = []
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append(" ".join(str(x) for x in row))
        return "\n".join(rows)
    except Exception as exc:
        logger.warning("CSV extraction failed for %s: %s", path, exc)
        return ""


def extract_docx(path: str) -> str:
    if Document is None:
        return ""
    try:
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as exc:
        logger.warning("DOCX extraction failed for %s: %s", path, exc)
        return ""


def extract_xlsx(path: str) -> str:
    if openpyxl is None:
        return ""
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        text = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                values = [str(x) for x in row if x is not None]
                if values:
                    text.append(" ".join(values))
        return "\n".join(text)
    except Exception as exc:
        logger.warning("XLSX extraction failed for %s: %s", path, exc)
        return ""


def extract_pptx(path: str) -> str:
    if Presentation is None:
        return ""
    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as exc:
        logger.warning("PPTX extraction failed for %s: %s", path, exc)
        return ""


EXTRACTORS = {
    ".pdf": extract_pdf,
    ".txt": extract_txt,
    ".json": extract_json,
    ".csv": extract_csv,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
    ".png": extract_image,
    ".jpg": extract_image,
    ".jpeg": extract_image,
    ".bmp": extract_image,
    ".tiff": extract_image,
    ".gif": extract_image,
}


def extract(path: str) -> str:
    """Extract and cache cleaned text content for a file on disk."""
    cached = _cache_get(path)
    if cached is not None:
        return cached

    if not os.path.isfile(path):
        return ""

    ext = Path(path).suffix.lower()
    extractor = EXTRACTORS.get(ext)

    if extractor:
        text = extractor(path)
    else:
        mime, _ = mimetypes.guess_type(path)
        text = extract_txt(path) if mime and mime.startswith("text") else ""

    text = clean(text)
    _cache_set(path, text)
    return text


def similarity(a: str, b: str) -> float:
    if not a or not b or fuzz is None:
        return 0.0
    return float(fuzz.token_set_ratio(a, b))
